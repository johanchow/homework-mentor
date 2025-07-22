"""
文件文字提取服务
支持从HTTP URL提取图片、PDF、PPT等文件中的文字
"""

import io
import re
import logging
from typing import Optional, List
from urllib.parse import urlparse
import easyocr
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
import docx
from utils.helpers import download_bytes_from_url

logger = logging.getLogger(__name__)

# 初始化EasyOCR读取器（支持中文和英文）
reader = easyocr.Reader(['ch_sim', 'en'])


class TextExtractor:
    """文字提取器"""
    
    def __init__(self):
        self.watermark_patterns = [
            r'水印|watermark|©|版权|版权所有',
            r'第\s*\d+\s*页|page\s*\d+',
            r'页码|page\s*number',
            r'www\.|http[s]?://',
            r'[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}',  # 邮箱
            r'\d{4}-\d{2}-\d{2}',  # 日期格式
            r'\d{2}:\d{2}:\d{2}',  # 时间格式
        ]
    
    def extract_text_from_url(self, file_url: str) -> str:
        """
        从HTTP URL提取文件中的文字
        
        Args:
            file_url: 文件的HTTP URL
            
        Returns:
            提取的文字字符串
        """
        try:
            # 获取文件字节流
            file_bytes = download_bytes_from_url(file_url)
            
            # 根据文件扩展名判断文件类型
            file_extension = self._get_file_extension(file_url)
            
            # 根据文件类型提取文字
            if self._is_image_file(file_extension):
                text = self._extract_text_from_image(file_bytes)
            elif file_extension.lower() == '.pdf':
                text = self._extract_text_from_pdf(file_bytes)
            elif file_extension.lower() in ['.ppt', '.pptx']:
                if file_extension.lower() == '.ppt':
                    file_bytes = self._convert_ppt_to_pptx(file_bytes)
                text = self._extract_text_from_ppt(file_bytes)
            elif file_extension.lower() in ['.doc', '.docx']:
                text = self._extract_text_from_word(file_bytes)
            else:
                raise ValueError(f"不支持的文件格式: {file_extension}")
            
            # 清理和过滤文字
            cleaned_text = self._clean_and_filter_text(text)
            
            return cleaned_text
            
        except Exception as e:
            logger.error(f"提取文件文字失败: {e}")
            raise
    
    def _get_file_extension(self, url: str) -> str:
        """获取文件扩展名"""
        parsed_url = urlparse(url)
        path = parsed_url.path
        return path[path.rfind('.'):] if '.' in path else ''
    
    def _is_image_file(self, extension: str) -> bool:
        """判断是否为图片文件"""
        image_extensions = ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.gif', '.webp']
        return extension.lower() in image_extensions
    
    def _extract_text_from_image(self, file_bytes: io.BytesIO) -> str:
        """从图片中提取文字（OCR）"""
        try:
            # 使用PIL打开图片
            image = Image.open(file_bytes)
            
            # 转换为numpy数组
            import numpy as np
            image_array = np.array(image)
            
            # 使用EasyOCR进行文字识别
            results = reader.readtext(image_array)
            
            # 提取文字
            texts = []
            for (bbox, text, confidence) in results:
                # 过滤低置信度的结果
                # if confidence > 0.5:
                texts.append(text)
            
            return '\n'.join(texts)
            
        except Exception as e:
            logger.error(f"图片OCR失败: {e}")
            raise
    
    def _extract_text_from_pdf(self, file_bytes: io.BytesIO) -> str:
        """从PDF中提取文字（包括图片中的文字）"""
        try:
            # 使用PyMuPDF打开PDF
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            
            texts = []
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # 1. 提取页面文字
                text = page.get_text()
                if text.strip():
                    texts.append(text)
                
                # 2. 提取页面中的图片并进行OCR
                image_texts = self._extract_text_from_pdf_images(page)
                if image_texts:
                    texts.extend(image_texts)
            
            doc.close()
            return '\n'.join(texts)
            
        except Exception as e:
            logger.error(f"PDF文字提取失败: {e}")
            raise
    
    def _extract_text_from_pdf_images(self, page) -> list:
        """从PDF页面中提取图片并进行OCR"""
        try:
            image_texts = []
            
            # 获取页面中的图片
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    # 获取图片数据
                    xref = img[0]
                    pix = fitz.Pixmap(page.parent, xref)
                    
                    # 转换为PIL Image
                    img_data = pix.tobytes("png")
                    img_bytes = io.BytesIO(img_data)
                    
                    # 进行OCR识别
                    ocr_text = self._extract_text_from_image(img_bytes)
                    if ocr_text.strip():
                        image_texts.append(ocr_text)
                    
                    pix = None  # 释放内存
                    
                except Exception as e:
                    logger.warning(f"PDF页面图片OCR失败 (图片{img_index + 1}): {e}")
                    continue
            
            return image_texts
            
        except Exception as e:
            logger.error(f"PDF图片提取失败: {e}")
            return []
    
    def _extract_text_from_ppt(self, file_bytes: io.BytesIO) -> str:
        """从PPT中提取文字（包括图片中的文字）"""
        try:
            # 使用python-pptx打开PPT
            prs = Presentation(file_bytes)
            
            texts = []
            for slide_index, slide in enumerate(prs.slides):
                slide_texts = []
                
                for shape in slide.shapes:
                    # 1. 提取文字形状中的文字
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                    
                    # 2. 提取图片中的文字
                    if shape.shape_type == 13:  # 图片类型
                        try:
                            image_text = self._extract_text_from_ppt_image(shape)
                            if image_text.strip():
                                slide_texts.append(f"[幻灯片{slide_index + 1}图片]: {image_text}")
                        except Exception as e:
                            logger.warning(f"PPT图片OCR失败 (幻灯片{slide_index + 1}): {e}")
                            continue
                
                if slide_texts:
                    texts.extend(slide_texts)
            
            return '\n'.join(texts)
            
        except Exception as e:
            logger.error(f"PPT文字提取失败: {e}")
            raise
    
    def _extract_text_from_ppt_image(self, shape) -> str:
        """从PPT图片形状中提取文字"""
        try:
            # 获取图片数据
            image = shape.image
            
            # 将图片数据转换为字节流
            img_bytes = io.BytesIO(image.blob)
            
            # 进行OCR识别
            return self._extract_text_from_image(img_bytes)
            
        except Exception as e:
            logger.error(f"PPT图片OCR失败: {e}")
            return ""
    
    def _extract_text_from_word(self, file_bytes: io.BytesIO) -> str:
        """从Word文档中提取文字（包括图片中的文字）"""
        try:
            # 使用python-docx打开Word文档
            doc = docx.Document(file_bytes)
            
            texts = []
            
            # 1. 提取段落文字
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    texts.append(paragraph.text)
            
            # 2. 提取表格文字
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text)
            
            # 3. 提取图片中的文字（如果存在）
            image_texts = self._extract_text_from_word_images(doc)
            if image_texts:
                texts.extend(image_texts)
            
            return '\n'.join(texts)
            
        except Exception as e:
            logger.error(f"Word文档文字提取失败: {e}")
            raise
    
    def _extract_text_from_word_images(self, doc) -> list:
        """从Word文档中提取图片并进行OCR"""
        try:
            image_texts = []
            
            # 遍历文档中的所有关系
            for rel in doc.part.rels.values():
                # 检查是否是图片关系
                if "image" in rel.target_ref:
                    try:
                        # 获取图片数据
                        image_data = rel.target_part.blob
                        img_bytes = io.BytesIO(image_data)
                        
                        # 进行OCR识别
                        ocr_text = self._extract_text_from_image(img_bytes)
                        if ocr_text.strip():
                            image_texts.append(f"[Word图片]: {ocr_text}")
                        
                    except Exception as e:
                        logger.warning(f"Word图片OCR失败: {e}")
                        continue
            
            return image_texts
            
        except Exception as e:
            logger.error(f"Word图片提取失败: {e}")
            return []
    
    def _clean_and_filter_text(self, text: str) -> str:
        """
        清理和过滤文字
        
        Args:
            text: 原始文字
            
        Returns:
            清理后的文字
        """
        if not text:
            return ""
        
        # 1. 清理多余的空格（保留换行符）
        # 将多个空格替换为单个空格，但保留换行符
        text = re.sub(r'[ \t]+', ' ', text)
        # 清理行首行尾的空格
        text = re.sub(r'^[ \t]+', '', text, flags=re.MULTILINE)
        text = re.sub(r'[ \t]+$', '', text, flags=re.MULTILINE)
        # 将多个连续换行符替换为单个换行符
        text = re.sub(r'\n\s*\n', '\n', text)
        
        # 2. 过滤水印、页码等无效文字
        for pattern in self.watermark_patterns:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE)
        
        # 3. 移除过短的行（可能是噪声）
        lines = text.split('\n')
        filtered_lines = []
        for line in lines:
            line = line.strip()
            if len(line) > 2:  # 保留长度大于2的行
                filtered_lines.append(line)
        
        # 4. 移除重复的行
        seen_lines = set()
        unique_lines = []
        for line in filtered_lines:
            if line not in seen_lines:
                seen_lines.add(line)
                unique_lines.append(line)
        
        # 5. 最终清理
        result = '\n'.join(unique_lines)
        result = result.strip()
        
        return result

    def _convert_ppt_to_pptx(self, file_bytes: io.BytesIO) -> io.BytesIO:
        """将PPT文件转换为PPTX文件"""
        pass

# 创建全局实例
text_extractor = TextExtractor()


def extract_text_from_file_url(file_url: str) -> str:
    """
    从文件URL提取文字的主函数
    
    Args:
        file_url: 文件的HTTP URL
        
    Returns:
        提取的文字字符串
    """
    return text_extractor.extract_text_from_url(file_url)


# 使用示例
if __name__ == "__main__":
    # 测试示例
    test_urls = [
        "https://clothing-try-on-1306401232.cos.ap-guangzhou.myqcloud.com/homework-mentor/1752582737-1b_13.pdf",
        "https://clothing-try-on-1306401232.cos.ap-guangzhou.myqcloud.com/homework-mentor/1752970995-11111.jpeg",
        # "https://clothing-try-on-1306401232.cos.ap-guangzhou.myqcloud.com/homework-mentor/1b_16.ppt",
        "https://clothing-try-on-1306401232.cos.ap-guangzhou.myqcloud.com/homework-mentor/1b_17.pptx",
    ]
    
    for url in test_urls:
        try:
            print(f"提取文件: {url}")
            text = extract_text_from_file_url(url)
            print(f"提取的文字: {text}")  # 只显示前200个字符
        except Exception as e:
            print(f"提取失败: {e}")
        print("-" * 50)
